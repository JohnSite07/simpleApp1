import boto3
import json
import os
from django.conf import settings
from botocore.exceptions import ClientError, NoCredentialsError
import logging

logger = logging.getLogger(__name__)

class AWSBedrockClient:
    def __init__(self):
        self.region = settings.AWS_REGION
        self.model_id = settings.AWS_BEDROCK_MODEL
        self.client = None
        self._initialize_client()
    
    def _initialize_client(self):
        """Initialize the Bedrock client with credentials from environment variables/settings"""
        try:
            self.client = boto3.client(
                'bedrock-runtime',
                region_name=self.region,
                aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
                aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
            )
            logger.info("AWS Bedrock client initialized successfully")
        except Exception as e:
            logger.error(f"Failed to initialize AWS Bedrock client: {e}")
            raise
    
    def generate_questions(self, text_content, quiz_type, num_questions):
        """Generate questions using AWS Bedrock"""
        try:
            prompt = self._create_prompt(text_content, quiz_type, num_questions)
            
            # Use the correct format for Claude 3
            if 'claude' in self.model_id.lower():
                response = self.client.invoke_model(
                    modelId=self.model_id,
                    body=json.dumps({
                        'anthropic_version': 'bedrock-2023-05-31',
                        'max_tokens': 4000,
                        'temperature': 0.7,
                        'top_p': 0.9,
                        'messages': [
                            {
                                'role': 'user',
                                'content': prompt
                            }
                        ]
                    })
                )
                
                response_body = json.loads(response['body'].read())
                generated_text = response_body.get('content', [{}])[0].get('text', '')
            else:
                # Fallback for other models
                response = self.client.invoke_model(
                    modelId=self.model_id,
                    body=json.dumps({
                        'prompt': prompt,
                        'max_tokens': 4000,
                        'temperature': 0.7,
                        'top_p': 0.9,
                        'stop_sequences': ['\n\nHuman:', '\n\nAssistant:']
                    })
                )
                
                response_body = json.loads(response['body'].read())
                generated_text = response_body.get('completion', '')
            
            return self._parse_questions(generated_text, quiz_type)
            
        except ClientError as e:
            logger.error(f"AWS Bedrock API error: {e}")
            raise Exception(f"Failed to generate questions: {e}")
        except Exception as e:
            logger.error(f"Error generating questions: {e}")
            raise
    
    def _create_prompt(self, text_content, quiz_type, num_questions):
        """Create a prompt for question generation"""
        if quiz_type == 'flashcards':
            prompt_template = f"""You are an expert educator creating flashcards from educational content. Based on the following text, create exactly {num_questions} high-quality flashcards.

Text content:
{text_content[:8000]}

Instructions:
1. Create exactly {num_questions} flashcards
2. Each flashcard should have a clear question on the front and a concise answer on the back
3. Focus on key concepts, definitions, important facts, and main ideas
4. Make questions that test understanding, not just memorization
5. Keep answers concise but complete
6. Ensure questions are relevant to the provided content

CRITICAL: You must respond with ONLY valid JSON. Do not include any explanatory text, markdown formatting, or code blocks. Start your response with {{ and end with }}.

Required JSON format:
{{
    "flashcards": [
        {{
            "question": "What is the main topic of this section?",
            "answer": "The main topic is..."
        }},
        {{
            "question": "Define the key concept mentioned.",
            "answer": "The key concept is..."
        }}
    ]
}}

Respond with ONLY the JSON object:"""
        else:  # multiple_choice
            prompt_template = f"""You are an expert educator creating multiple choice questions from educational content. Based on the following text, create exactly {num_questions} high-quality multiple choice questions.

Text content:
{text_content[:8000]}

Instructions:
1. Create exactly {num_questions} multiple choice questions
2. Each question should have 4 options (A, B, C, D)
3. Only one option should be correct
4. Focus on key concepts, definitions, important facts, and main ideas
5. Make questions that test understanding, not just memorization
6. Ensure all distractors (wrong answers) are plausible
7. Ensure questions are relevant to the provided content

CRITICAL: You must respond with ONLY valid JSON. Do not include any explanatory text, markdown formatting, or code blocks. Start your response with {{ and end with }}.

Required JSON format:
{{
    "questions": [
        {{
            "question": "What is the main topic discussed?",
            "options": {{
                "A": "Option A text",
                "B": "Option B text", 
                "C": "Option C text",
                "D": "Option D text"
            }},
            "correct_answer": "A"
        }}
    ]
}}

Respond with ONLY the JSON object:"""
        return prompt_template
    
    def _parse_questions(self, generated_text, quiz_type):
        """Parse the generated text into structured question data"""
        try:
            # Clean the response text
            generated_text = generated_text.strip()
            logger.debug(f"Raw generated text: {generated_text[:500]}...")
            
            # Try to find JSON in the response
            json_start = generated_text.find('{')
            json_end = generated_text.rfind('}') + 1
            
            if json_start == -1 or json_end == 0:
                # If no JSON found, try to extract from markdown code blocks
                if '```json' in generated_text:
                    start = generated_text.find('```json') + 7
                    end = generated_text.find('```', start)
                    if end != -1:
                        json_str = generated_text[start:end].strip()
                    else:
                        raise ValueError("No JSON found in response")
                elif '```' in generated_text:
                    # Try to extract from any code block
                    start = generated_text.find('```') + 3
                    end = generated_text.find('```', start)
                    if end != -1:
                        json_str = generated_text[start:end].strip()
                    else:
                        raise ValueError("No JSON found in response")
                else:
                    # Try to extract JSON-like content
                    lines = generated_text.split('\n')
                    json_lines = []
                    in_json = False
                    for line in lines:
                        if '{' in line or in_json:
                            in_json = True
                            json_lines.append(line)
                        if '}' in line and in_json:
                            break
                    if json_lines:
                        json_str = '\n'.join(json_lines)
                    else:
                        raise ValueError("No JSON found in response")
            else:
                json_str = generated_text[json_start:json_end]
            
            logger.debug(f"Extracted JSON string: {json_str[:500]}...")
            
            # Parse JSON
            data = json.loads(json_str)
            
            if quiz_type == 'flashcards':
                flashcards = data.get('flashcards', [])
                if not flashcards:
                    # Try alternative key names
                    flashcards = data.get('questions', [])
                    if not flashcards:
                        flashcards = data.get('cards', [])
                
                # Validate flashcards structure
                valid_flashcards = []
                for i, flashcard in enumerate(flashcards):
                    if isinstance(flashcard, dict):
                        question = flashcard.get('question', flashcard.get('front', ''))
                        answer = flashcard.get('answer', flashcard.get('back', ''))
                        if question and answer:
                            valid_flashcards.append({
                                'question': question,
                                'answer': answer
                            })
                        else:
                            logger.warning(f"Flashcard {i} missing question or answer: {flashcard}")
                    else:
                        logger.warning(f"Flashcard {i} is not a dictionary: {flashcard}")
                
                if not valid_flashcards:
                    raise ValueError("No valid flashcards found in response")
                
                logger.info(f"Successfully parsed {len(valid_flashcards)} flashcards")
                return valid_flashcards
            else:
                questions = data.get('questions', [])
                # Validate questions structure
                valid_questions = []
                for i, question in enumerate(questions):
                    if isinstance(question, dict):
                        q_text = question.get('question', '')
                        options = question.get('options', {})
                        correct = question.get('correct_answer', '')
                        if q_text and options and correct:
                            valid_questions.append(question)
                        else:
                            logger.warning(f"Question {i} missing required fields: {question}")
                    else:
                        logger.warning(f"Question {i} is not a dictionary: {question}")
                
                if not valid_questions:
                    raise ValueError("No valid questions found in response")
                
                logger.info(f"Successfully parsed {len(valid_questions)} questions")
                return valid_questions
                
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON response: {e}")
            logger.error(f"Generated text: {generated_text}")
            raise Exception("Failed to parse AI response - invalid JSON format")
        except Exception as e:
            logger.error(f"Error parsing questions: {e}")
            logger.error(f"Generated text: {generated_text}")
            raise

class DocumentProcessor:
    """Handle document processing and text extraction"""
    
    @staticmethod
    def extract_text_from_file(file_path):
        """Extract text from various file types"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension in ['.txt', '.md']:
                return DocumentProcessor._extract_text_plain(file_path)
            elif file_extension in ['.pdf']:
                return DocumentProcessor._extract_text_pdf(file_path)
            elif file_extension in ['.docx', '.doc']:
                return DocumentProcessor._extract_text_docx(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return DocumentProcessor._extract_text_pptx(file_path)
            else:
                # For unsupported file types, try to extract as plain text
                return DocumentProcessor._extract_text_plain(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            raise
    
    @staticmethod
    def _extract_text_plain(file_path):
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
    
    @staticmethod
    def _extract_text_pdf(file_path):
        """Extract text from PDF files"""
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except ImportError:
            logger.warning("PyPDF2 not installed, treating PDF as plain text")
            return DocumentProcessor._extract_text_plain(file_path)
    
    @staticmethod
    def _extract_text_docx(file_path):
        """Extract text from DOCX files"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            logger.warning("python-docx not installed, treating DOCX as plain text")
            return DocumentProcessor._extract_text_plain(file_path)
    
    @staticmethod
    def _extract_text_pptx(file_path):
        """Extract text from PPTX files"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logger.warning("python-pptx not installed, treating PPTX as plain text")
            return DocumentProcessor._extract_text_plain(file_path) 